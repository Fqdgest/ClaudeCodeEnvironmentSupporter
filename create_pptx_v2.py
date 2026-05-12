"""
Generate professional PPTX from Word document content.
Layout: left-top title, left-text + right-diagram columns, bottom caption.
LaTeX formulas in clearly marked boxes for PowerPoint native conversion.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ============================================================
# Constants
# ============================================================
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# Colors
NAVY    = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BG = RGBColor(0x0F, 0x1B, 0x33)
TEAL    = RGBColor(0x08, 0x91, 0xB2)
CORAL   = RGBColor(0xF9, 0x61, 0x67)
GOLD    = RGBColor(0xF9, 0xE7, 0x95)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF0, 0xF3, 0xFA)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
MUTED   = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
CREAM   = RGBColor(0xF8, 0xF6, 0xF0)

FONT_TITLE = "Microsoft YaHei"
FONT_BODY  = "Microsoft YaHei"
FONT_MONO  = "Consolas"

# Margins
ML = Inches(0.55)   # left margin
MR = Inches(0.55)   # right margin
MT = Inches(0.28)   # top margin for title
COL_GAP = Inches(0.25)  # gap between left and right columns

# Column widths
LEFT_COL_W = Inches(5.2)
RIGHT_COL_W = Inches(3.6)
RIGHT_COL_X = ML + LEFT_COL_W + COL_GAP
CONTENT_TOP = Inches(1.05)
CONTENT_BODY_H = Inches(3.85)
BOTTOM_Y = Inches(5.05)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ============================================================
# Helpers
# ============================================================
def add_blank_slide(bg_color=WHITE):
    """Add a blank slide with optional background."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    if bg_color != WHITE:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return slide

def add_left_accent_bar(slide, color=NAVY):
    """Thin vertical accent bar on left edge."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.1), SLIDE_H)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_title(slide, text, y=MT):
    """Add slide title at top-left."""
    txBox = slide.shapes.add_textbox(ML, y, Inches(9.0), Inches(0.55))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_TITLE
    # underline accent
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, ML, y + Inches(0.52), Inches(0.7), Inches(0.04))
    shp.fill.solid()
    shp.fill.fore_color.rgb = TEAL
    shp.line.fill.background()
    return txBox

def add_left_text(slide, bullets, y=CONTENT_TOP, h=CONTENT_BODY_H, max_chars_per_bullet=120):
    """Add left-column text with formatted bullets."""
    txBox = slide.shapes.add_textbox(ML, y, LEFT_COL_W, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Label (bold)
        if "label" in item:
            run = p.add_run()
            run.text = item["label"]
            run.font.size = Pt(12)
            run.font.bold = True
            run.font.color.rgb = NAVY
            run.font.name = FONT_BODY

        # Description
        if "desc" in item:
            run = p.add_run()
            run.text = item["desc"][:max_chars_per_bullet]
            run.font.size = Pt(11)
            run.font.color.rgb = DARK_TEXT
            run.font.name = FONT_BODY

        # Sub bullets
        if "sub" in item:
            for sb in item["sub"]:
                p2 = tf.add_paragraph()
                p2.space_before = Pt(0)
                p2.space_after = Pt(0)
                run = p2.add_run()
                run.text = "    · " + sb[:max_chars_per_bullet]
                run.font.size = Pt(10)
                run.font.color.rgb = MUTED
                run.font.name = FONT_BODY

        p.space_after = Pt(4)
    return txBox

def add_right_diagram_area(slide, y=CONTENT_TOP, h=CONTENT_BODY_H):
    """Add right-column diagram area with two image placeholders and a caption box."""
    # Background rect for diagram area
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, RIGHT_COL_X, y, RIGHT_COL_W, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = OFF_WHITE
    bg.line.color.rgb = LIGHT_GRAY
    bg.line.width = Pt(0.5)
    bg.name = "diagram_bg"

    # Upper image placeholder
    img1_h = (h - Inches(0.5)) * 0.42
    img1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        RIGHT_COL_X + Inches(0.2), y + Inches(0.15),
        RIGHT_COL_W - Inches(0.4), img1_h)
    img1.fill.solid()
    img1.fill.fore_color.rgb = WHITE
    img1.line.color.rgb = LIGHT_GRAY
    img1.line.width = Pt(0.5)
    img1.name = "img_placeholder_1"
    # Label
    tf = img1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "📊 图表区域 1"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.font.name = FONT_BODY

    # Lower image placeholder
    img2_y = y + img1_h + Inches(0.35)
    img2_h = (h - Inches(0.5)) * 0.35
    img2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        RIGHT_COL_X + Inches(0.2), img2_y,
        RIGHT_COL_W - Inches(0.4), img2_h)
    img2.fill.solid()
    img2.fill.fore_color.rgb = WHITE
    img2.line.color.rgb = LIGHT_GRAY
    img2.line.width = Pt(0.5)
    img2.name = "img_placeholder_2"
    tf = img2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "📊 图表区域 2"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.font.name = FONT_BODY

    return img1, img2

def add_formula_box(slide, latex_text, y=None, x=None, w=None, label="LATEX"):
    """Add a formula box with LaTeX code that PPT can convert natively.
       User selects text → Insert Equation → paste LaTeX → converts to editable math.
    """
    if y is None:
        y = BOTTOM_Y
    if x is None:
        x = ML
    if w is None:
        w = Inches(9.0)

    # Small accent line above formula
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y - Inches(0.05), Inches(0.4), Inches(0.02))
    shp.fill.solid()
    shp.fill.fore_color.rgb = CORAL
    shp.line.fill.background()

    # Formula box
    txBox = slide.shapes.add_textbox(x, y + Inches(0.02), w, Inches(0.48))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{label}:  {latex_text}" if label else latex_text
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED
    run.font.name = FONT_MONO
    run.font.italic = True
    return txBox

def add_caption(slide, text, y=CONTENT_TOP + CONTENT_BODY_H + Inches(0.05)):
    """Add a small caption text below the main content area."""
    txBox = slide.shapes.add_textbox(ML, y, Inches(9.0), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.font.name = FONT_BODY
    p.font.italic = True
    return txBox

def make_content_slide(title, bullets, formulas=None, captions=None):
    """Create a standard content slide with left-text + right-diagram layout."""
    slide = add_blank_slide(WHITE)
    add_left_accent_bar(slide)
    add_title(slide, title)
    if captions:
        add_caption(slide, captions)
    add_left_text(slide, bullets)
    add_right_diagram_area(slide)
    if formulas:
        for f in formulas:
            add_formula_box(slide, f["latex"], y=f.get("y"), label=f.get("label", "LATEX"))
    return slide

def make_section_slide(number, title, subtitle=""):
    """Create a dark section divider slide."""
    slide = add_blank_slide(DARK_BG)
    # Large number
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(2), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.font.name = FONT_TITLE
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(8.5), Inches(1.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE
    # Accent line
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(3.2), Inches(1.5), Inches(0.05))
    shp.fill.solid()
    shp.fill.fore_color.rgb = CORAL
    shp.line.fill.background()
    # Subtitle
    if subtitle:
        txBox = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(8.5), Inches(0.6))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = MUTED
        p.font.name = FONT_BODY
    return slide

def make_title_slide():
    """Create the main title slide."""
    slide = add_blank_slide(DARK_BG)
    # Top decorative line
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.0), Inches(1.8), Inches(0.06))
    shp.fill.solid()
    shp.fill.fore_color.rgb = TEAL
    shp.line.fill.background()
    # Main title
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "串联支座隔震结构数值模型"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT_TITLE
    p2 = tf.add_paragraph()
    p2.text = "9自由度理论框架深度审查报告"
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = TEAL
    p2.font.name = FONT_TITLE
    # Subtitle line
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(3.0), Inches(8.6), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Koh-Kelly 两弹簧离散模型 · 变分原理推导 · 舒尔补刚度凝聚 · 非线性动力时程分析"
    p.font.size = Pt(13)
    p.font.color.rgb = MUTED
    p.font.name = FONT_BODY
    # Bottom line
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(3.8), Inches(2.0), Inches(0.04))
    shp.fill.solid()
    shp.fill.fore_color.rgb = CORAL
    shp.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(8.6), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "iFLOW 结构抗震分析团队  |  2026年5月"
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED
    p.font.name = FONT_BODY
    return slide

# ============================================================
# BUILD SLIDES
# ============================================================

# ---- SLIDE 1: Title ----
make_title_slide()

# ---- SLIDE 2: Agenda ----
make_section_slide("01", "报告结构总览", "本文档涵盖两大篇章，共六个核心技术板块")

# ---- SLIDE 3: Agenda detail ----
s = add_blank_slide(WHITE)
add_left_accent_bar(s)
add_title(s, "报告内容导航")

agenda = [
    {"label": "第一篇 ", "desc": "串联支座隔震结构数值模型及9自由度理论框架深度审查"},
    {"label": "    第1章 ", "desc": "绪论与串联隔震体系技术背景"},
    {"label": "    第2章 ", "desc": "隔震单元大变形理论演进 (3DOF→5DOF体系)"},
    {"label": "    第3章 ", "desc": "五自由度单支座变分推导与刚度矩阵解析"},
    {"label": "    第4章 ", "desc": "9自由度串联超单元拓扑架构与状态确定"},
    {"label": "第二篇 ", "desc": "串联支座隔震结构非线性动力学数值模型 (深度扩展版)"},
    {"label": "    第5章 ", "desc": "整体建模思路与坐标系修正 · 非线性刚度 · 单/串联支座数学模型"},
    {"label": "    第6章 ", "desc": "舒尔补刚度凝聚校核 · 上部多质点体系全局耦合"},
]
add_left_text(s, agenda, h=Inches(4.3))
add_right_diagram_area(s)

# ---- SLIDE 4: Chapter 1 intro ----
make_section_slide("02", "第1章：绪论与串联隔震体系技术背景",
                   "从传统基底隔震到串联隔震体系的技术演进")

# ---- SLIDE 5: Background ----
make_content_slide("传统单层基底隔震的挑战与串联体系的诞生", [
    {"label": "传统基座隔震的三大局限", "desc": ""},
    {"label": "位移超限", "desc": "：近断层长周期速度脉冲下，隔震层位移远超设计容许值"},
    {"label": "支座受拉屈曲", "desc": "：强竖向地震动分量引发支座轴向拉力，橡胶层-钢板层离"},
    {"label": "三维耦合失稳", "desc": "：平动-转动惯性耦合导致动力响应难以预测和控制"},
    {"label": "串联隔震体系 (Series Isolation System)", "desc": ""},
    {"label": "核心思想", "desc": "：叠层橡胶支座 + 摩擦摆/金属阻尼器通过刚性中间连接板串联拼装"},
    {"label": "机制", "desc": "：多级物理界面协同变形 → 提升变形能力与耗散裕度"},
    {"label": "理论基础", "desc": "：Koh-Kelly (1987) 两弹簧离散宏观力学模型 → 预测大剪切稳定性和刚度退化"},
], formulas=[{"latex": r"$K_h(s)$ = $K_{h0}$ \cdot (1 - \tanh(\alpha \cdot |s|/$T_r$))$", "label": "Koh-Kelly剪切刚度退化"}])

# ---- SLIDE 6: Boundary condition problem ----
make_content_slide("串联体系中边界条件的拓扑重构", [
    {"label": "经典Koh-Kelly模型的边界假设", "desc": ""},
    {"label": "3DOF模型", "desc": "：底板固结于基础，顶板与上部结构刚性固结，无独立转角 ($\\theta_v$=$\\theta_b$=0)"},
    {"label": "串联体系的拓扑变化", "desc": ""},
    {"label": "上层支座底部", "desc": "：不再固结于基础，而是连接于具备水平平动+竖向平动+绝对转动自由度的中间连接板"},
    {"label": "传统模型的致命缺陷", "desc": ""},
    {"label": "", "desc": "若沿用顶部固结的3DOF/4DOF模型，无法满足中间板与上下支座的复杂变形协调，也不能真实映射外力剪力、动态轴力与支座弯矩共同施加于底板的绝对空间弯矩平衡"},
    {"label": "解决方案", "desc": "：构建5-DOF单支座微观单元 → 9-DOF串联超单元数值模型"},
], formulas=[{"latex": r"v^{(5DOF)} = [s_x, $s_z$, \theta_h, \theta_v, \theta_b]^T", "label": "5-DOF状态向量"}])

# ---- SLIDE 7: Chapter 2 intro ----
make_section_slide("03", "第2章：隔震单元大变形理论的演进",
                   "从基准3自由度到5自由度体系 · 均分原则 · 物理对称性重构")

# ---- SLIDE 8: 3DOF limitations ----
make_content_slide("传统3-DOF Koh-Kelly模型的力学本质与局限", [
    {"label": "3-DOF内部自由度", "desc": "：水平剪切变形 $s_x$、竖向压缩 $s_z$、支座主轴倾角 θ"},
    {"label": "核心简化假定", "desc": "：顶板与底板均无独立刚体转角"},
    {"label": "力学后果", "desc": ""},
    {"label": "静力凝聚", "desc": "：顶部和底部两个物理旋转弹簧被数学等效为并联的单一等效旋转刚度 $K_\\theta$"},
    {"label": "信息丢失", "desc": "：掩盖了支座内部弯曲应变能沿高度方向的对称分布规律"},
    {"label": "适用边界", "desc": "：仅在单层基础隔震（仅平动剪切）时具备一定精度"},
    {"label": "串联场景下的断裂", "desc": "：无法输出端部局部弯矩，物理上切断了与相邻构件的转动耦合"},
], formulas=[{"latex": r"v^{(3DOF)} = [s_x, $s_z$, \theta]^T, \quad K_\theta = K_{\theta}^{top} + K_{\theta}^{bot}", "label": "经典3-DOF"}])

# ---- SLIDE 9: Equal distribution ----
make_content_slide("均分原则 (Equal Distribution Principle)", [
    {"label": "物理基础", "desc": "：叠层橡胶支座在微观材料构造与宏观几何上完全上下对称"},
    {"label": "均分原则", "desc": ""},
    {"label": "", "desc": "将原3DOF模型中打包的总抗弯刚度 $K_\\theta$0 严格均分为两部分"},
    {"label": "", "desc": "顶部弹簧与底部弹簧各自获得 0.5 × $K_\\theta$0"},
    {"label": "剪切非线性退化", "desc": ""},
    {"label": "P-Δ效应", "desc": "：大变形下轴向压力改变橡胶层应力分布，抗弯能力随水平剪切位移增加而退化"},
    {"label": "统一本构", "desc": "：顶/底旋转弹簧均赋予退化参数 $C_\\theta$，即 $K_\\theta$(s) = 0.5·$K_\\theta$0·(1-$C_\\theta$·|s_x|)"},
    {"label": "向下兼容性验证", "desc": ""},
    {"label": "", "desc": "锁死 $\\theta_v$=0$ 且 $\\theta_b$=0$ → 两等刚度弹簧自动重组为并联 → 精准退化回经典3DOF模型"},
], formulas=[{"latex": r"K_{\theta}^{top} = K_{\theta}^{bot} = \frac{1}{2} K_{\theta 0} \cdot (1 - C_\theta \cdot |s_x|)", "label": "均分退化公式"}])

# ---- SLIDE 10: Model comparison table ----
s = add_blank_slide(WHITE)
add_left_accent_bar(s)
add_title(s, "不同维度 Koh-Kelly 模型的演进对比")

# Table
rows, cols = 6, 4
tbl = s.shapes.add_table(rows, cols, ML, CONTENT_TOP, Inches(9.0), Inches(3.2)).table
tbl.columns[0].width = Inches(1.5)
tbl.columns[1].width = Inches(2.8)
tbl.columns[2].width = Inches(2.5)
tbl.columns[3].width = Inches(2.2)
headers = ["模型维度", "内部状态向量构成", "边界假定", "抗弯刚度分配逻辑"]
data = [
    ["经典 3-DOF", "v=[s_x, $s_z$, θ]^T", "顶板$\theta_v$=0, 底板$\theta_b$=0", "集中并联为单一总刚度"],
    ["改进 4-DOF", "v=[s_x, $s_z$, θ_h, θ_v]^T", "顶板独立转角, 底板$\theta_b$=0", "均分+剪切退化"],
    ["通用 5-DOF", "v=[s_x, $s_z$, θ_h, θ_v, θ_b]^T", "顶底板均独立转角", "顶底完全对称非线性"],
]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = FONT_BODY
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
for i, row in enumerate(data):
    for j, val in enumerate(row):
        cell = tbl.cell(i+1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_TEXT
            p.font.name = FONT_BODY
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = OFF_WHITE

add_formula_box(s, r"v^{(5DOF)} = [s_x,\; $s_z$,\; \theta_h,\; \theta_v,\; \theta_b]^T", y=Inches(4.4), label="通用5-DOF")
add_right_diagram_area(s)

# ---- SLIDE 11: Chapter 3 intro ----
make_section_slide("04", "第3章：五自由度单支座变分推导",
                   "变应变能泛函 · 内力向量 · 材料刚度矩阵 · 几何刚度截断缺陷")

# ---- SLIDE 12: Variational derivation ----
make_content_slide("基于变分原理的能量方程推导框架", [
    {"label": "方法论优势", "desc": "：采用变分原理(Variational Principle)推导，避免传统牛顿力学在复杂自由度系统中的漏项问题"},
    {"label": "总变应变能 Π", "desc": "：由五部分机械能组成"},
    {"label": "① 水平剪切能", "desc": "：½·$K_h$·s_x²"},
    {"label": "② 竖向压缩能", "desc": "：$\\frac{1}{2}K_v$·s_z²"},
    {"label": "③④ 顶/底相对弯曲能", "desc": "：½·$K_\\theta$t·(θ_h-θ_v)² + ½·$K_\\theta$b·(θ_v-θ_b)²  (对称)"},
    {"label": "⑤ 形心轴弯曲能", "desc": "：支座主轴倾角变化产生的应变能贡献"},
    {"label": "内力向量", "desc": "：F_int = ∂Π/∂v_i  (链式法则求偏导)"},
], formulas=[{"latex": r"\Pi = \frac{1}{2}$K_h$ $s_x^2$ + \frac{1}{2}$K_v$ $s_z^2$ + \frac{1}{2}K_{\theta t}(\theta_h-\theta_v)^2 + \frac{1}{2}K_{\theta b}(\theta_v-\theta_b)^2 + \cdots", "label": "变应变能泛函"}])

# ---- SLIDE 13: Shear deterioration ----
make_content_slide("剪切劣化效应 (Shear Deterioration Effect)", [
    {"label": "物理现象", "desc": "：大剪切变形时橡胶层-钢板受力状态改变 → 抗弯刚度 $K_\\theta$ 下降"},
    {"label": "能量守恒视角", "desc": ""},
    {"label": "退化", "desc": "：∂$K_\\theta$/∂$s_x$ < 0 (导数为负值)"},
    {"label": "释放", "desc": "：弯曲刚度退化释放的应变能转化为沿剪切方向的附加负剪力"},
    {"label": "数学机制", "desc": ""},
    {"label": "交叉导数项", "desc": "：F_shear 表达式中的 ∂$K_\\theta$/∂$s_x$ 项——弯曲变形与剪切承载力之间的交叉反馈"},
    {"label": "宏观表现", "desc": "：强轴压+大转角下支座剪切失稳加速"},
    {"label": "工业意义", "desc": "：若未正确建模此项，将低估大变形下的支座剪切失效风险"},
], formulas=[{"latex": r"$F_{shear}$ = \frac{\partial\Pi}{\partial $s_x$} - \sum_i \frac{\partial\Pi}{\partial\theta_i} \cdot \frac{\partial\theta_i}{\partial $s_x$}, \quad \frac{\partial K_\theta}{\partial $s_x$} < 0", "label": "剪切劣化交叉导数"}])

# ---- SLIDE 14: Material stiffness ----
make_content_slide("材料刚度矩阵 $K_{mat}$ 的完备性与对称性", [
    {"label": "定义", "desc": "：$K_{mat}$ = $\partial F_{int}/\partial v$ = $\partial^2\Pi/\partial v^2$ —— 变应变能的二阶偏导数"},
    {"label": "数学性质", "desc": ""},
    {"label": "", "desc": "完美的 5×5 对称矩阵"},
    {"label": "", "desc": "所有非主对角线耦合元素保持严格对称性 ($K_{mat}$)_{ij} = ($K_{mat}$)_{ji}"},
    {"label": "", "desc": "高阶导数的完备性确保材料本构层面不发生病态非对称"},
    {"label": "关键耦合项示例", "desc": ""},
    {"label": "∂²Π/(∂s_x·∂θ_v)", "desc": "：精确量化底板转角偏差对初始剪切刚度的削弱矩阵元"},
    {"label": "∂²Π/(∂θ_h·∂θ_v)", "desc": "：顶板-形心轴转动耦合刚度"},
], formulas=[{"latex": r"$K_{mat}$ = \frac{\partial^2\Pi}{\partial v^2} \in \mathbb{R}^{5\times 5}, \quad $K_{mat}$ = $K_{mat}$^T \quad \text{(严格对称)}", "label": "材料刚度矩阵"}])

# ---- SLIDE 15: Geometric stiffness truncation ----
make_content_slide("几何刚度矩阵 $K_{geo}$ 的截断缺陷分析", [
    {"label": "几何刚度的来源", "desc": "：外部节点力在变形位形下的雅可比偏导——大变形坐标投影修正"},
    {"label": "发现的截断缺陷", "desc": ""},
    {"label": "", "desc": "非零元素仅被限定于竖向压缩 $s_z$ 和形心轴倾角 θ_h 相关的局部子块内"},
    {"label": "", "desc": "完全丢弃了 θ_v 和 θ_b 对位移雅可比矩阵 T 的二阶修正"},
    {"label": "物理后果", "desc": ""},
    {"label": "剧烈摇摆工况", "desc": "：中间板转角与顶板转角产生巨大绝对偏差"},
    {"label": "端部转角", "desc": "：引发橡胶层二次局部曲率变形 → 改变轴力 P 在剪切平面上的真实投影"},
    {"label": "临界屈曲", "desc": "：切线刚度评估偏乐观 → Newton-Raphson二次收敛特性被破坏 → 伪收敛或迭代发散"},
], formulas=[{"latex": r"$K_{tangent}$ = $K_{mat}$ + $K_{geo}$, \quad $K_{geo}$ = \frac{\partial}{\partial v}($F_{ext}$ \circ \Phi), \quad \text{截断: } \frac{\partial T}{\partial\theta_v} \approx 0", "label": "几何刚度截断"}])

# ---- SLIDE 16: Chapter 4 intro ----
make_section_slide("05", "第4章：9自由度串联超单元",
                   "拓扑架构 · 9维状态向量 · 闭环残差方程 · 舒尔补刚度凝聚")

# ---- SLIDE 17: Super-element topology ----
make_content_slide("9-DOF串联超单元的物理拓扑", [
    {"label": "超单元方法 (Super-Element)", "desc": "：在全局组装阶段向主求解器屏蔽复杂的内部动态自由度，仅暴露宏观边界自由度"},
    {"label": "拓扑构成", "desc": ""},
    {"label": "下支座 (单元1)", "desc": "：底板与基础固结 ($U_{bot}$=0)，顶板连接中间板"},
    {"label": "上支座 (单元2)", "desc": "：底部连接中间板，顶部连接刚性底板 (由全局响应驱动)"},
    {"label": "中间连接板", "desc": "：刚性自由体——承载水平剪力、竖向轴力的传力枢纽"},
    {"label": "内外自由度划分", "desc": ""},
    {"label": "内部 (9 DOF)", "desc": "：上下支座内部变形 + 中间板绝对空间位姿"},
    {"label": "外部 (3 DOF)", "desc": "：仅刚性底板的 $U_{top}$=[$X_{top}$, $Z_{top}$, Θ_top]^T 与主结构交互"},
], formulas=[{"latex": r"v_{9DOF} = [$s_{x1}$, $s_{z1}$, \theta_{h1}, \theta_{v1}, X_{mid}, Z_{mid}, \Theta_{mid}, $s_{x2}$, $s_{z2}$, \theta_{h2}, \theta_{v2}]^T", "label": "9维状态向量"}])

# ---- SLIDE 18: Residual equations I ----
make_content_slide("闭环残差方程组 (I)：下支座行为", [
    {"label": "第一部分方程 R₁(s_x1)→R₃(θ_h1)", "desc": ""},
    {"label": "", "desc": "将下支座内部剪切、压缩变形及主轴转角，通过正余弦大变形映射"},
    {"label": "", "desc": "与中间连接板绝对平移空间位置 [$X_{mid}$, $Z_{mid}$] 直接绑定"},
    {"label": "运动学协调", "desc": ""},
    {"label": "s_x1", "desc": "：剪切变形 ⇄ 中间板水平平动 + 支座大转动几何投影"},
    {"label": "s_z1", "desc": "：压缩变形 ⇄ 中间板竖向沉降 + 主轴倾角几何耦合"},
    {"label": "力矩平衡 R₃", "desc": "：下支座主轴倾覆力矩 = 回转刚度约束 —— 中间板转角 Θ_mid 作为下支座顶部边界条件"},
], formulas=[{"latex": r"R_1: $s_{x1}$ - f_{geom}^{(1)}(X_{mid}, Z_{mid}, \Theta_{mid}) = 0", "label": "下支座残差"}])

# ---- SLIDE 19: Residual equations II ----
make_content_slide("闭环残差方程组 (II)：上支座 + 中间板连续性", [
    {"label": "第二部分 R₄→R₆：上支座网络", "desc": ""},
    {"label": "边界特点", "desc": "：两端边界均未知——全局驱动顶板 $U_{top}$ + 自由浮动中间板"},
    {"label": "剪切变形", "desc": "：s_x2 由 $U_{top}$ 与中间板相对位移控制"},
    {"label": "弯矩平衡", "desc": "：含 P-Δ 驱动项 + $K_\\theta$ 旋转刚度组合恢复力矩"},
    {"label": "第三部分 R₇→R₉：中间板静力连续性", "desc": ""},
    {"label": "R₇(R₈)", "desc": "：保障水平剪力/竖向轴力穿过串联体系的绝对传递连续性"},
    {"label": "R₉", "desc": "：核心——绝对弯矩连续+碰撞非线性锁止：$M_b$ot + M_top + M_contact = 0"},
], formulas=[{"latex": r"R_9: $M_{bot}$ + $M_{top}$ + $M_{contact}$(\Theta_{mid}) = 0, \quad $M_{contact}$ = \text{Penalty}\cdot\langle|\Theta_{mid}|-\Theta_{lim}\rangle", "label": "中间板弯矩连续"}])

# ---- SLIDE 20: Schur complement ----
make_content_slide("舒尔补 (Schur Complement) 切线刚度凝聚", [
    {"label": "为什么需要舒尔补？", "desc": ""},
    {"label": "", "desc": "中间板大转角、P-Δ局部负刚度、接触刚度突跳 → 直接进全局求解器极易崩溃"},
    {"label": "数学原理", "desc": ""},
    {"label": "混合列式", "desc": "：构建拉格朗日鞍点系统 [K_int  Tᵀ; T  0]·[dv; dF] = [0; dU]"},
    {"label": "舒尔补", "desc": "：对左上角 K_int 求舒尔补消去微观变分 dv → 获得宏观刚度"},
    {"label": "凝聚公式", "desc": "：K_global = (T·K_int⁻¹·Tᵀ)⁻¹"},
    {"label": "工程价值", "desc": ""},
    {"label": "", "desc": "无论中间节点经历何种严重力学退化，输出的宏观切线刚度始终平滑一致"},
    {"label": "", "desc": "所有转角碰撞惩罚与二阶几何不稳定性在「黑盒」中被消化并封闭"},
], formulas=[{"latex": r"$K_{global}$ = \left( T \cdot $K_{int}$^{-1} \cdot T^T \right)^{-1}, \quad $K_{int}$ = \frac{\partial R_{9DOF}}{\partial v_{9DOF}}", "label": "舒尔补凝聚公式"}])

# ---- SLIDE 21: Part 2 intro ----
make_section_slide("06", "第二篇：非线性动力学数值模型 (深度扩展版)",
                   "绝对坐标系 · 符号自洽性 · 非线性刚度 · 单/串联支座 · 全局耦合")

# ---- SLIDE 22: Absolute coordinate system ----
make_content_slide("绝对坐标系的革新与系统解耦", [
    {"label": "传统相对坐标系的痛点", "desc": ""},
    {"label": "质量矩阵", "desc": "：M 变为满阵(平动-转动惯性耦合) → 求逆代价大 → 截断误差"},
    {"label": "多维地震动", "desc": "：相对位移在基础摇摆转角时产生极其复杂的惯性耦合项"},
    {"label": "绝对坐标系的优势", "desc": ""},
    {"label": "对角化", "desc": "：M = diag(m₁, m₂, ..., m_N) — 完美对角矩阵"},
    {"label": "解耦", "desc": "：运动学映射矩阵退化为纯单位阵 → 行波效应/多维激励直观严密"},
    {"label": "嵌套求解架构", "desc": ""},
    {"label": "Newmark-β", "desc": "：全局无条件稳定时间步进 → 刚性底板下发试探位移 → 支座Newton-Raphson内迭代 → 刚度凝聚上传"},
], formulas=[{"latex": r"M = \operatorname{diag}(m_1, m_2, \ldots, m_N), \quad T_{kin} = I_{total\_dof}", "label": "绝对坐标质量矩阵"}])

# ---- SLIDE 23: Symbol sign correction ----
make_content_slide("⚠️ 关键修正：竖向变形 $s_z$ 与轴力 P 的符号自洽性", [
    {"label": "物理悖论", "desc": ""},
    {"label": "", "desc": "支座受压：$s_z$ < 0, P < 0 (向上为正坐标系)"},
    {"label": "错误写法", "desc": "：$h_{eff}$ = h₀ - $s_z$ → 受压时高度反而增加 (计算结果 h₀ + |s_z|)！"},
    {"label": "严格修正：$h_{eff}$ = h₀ + $s_z$", "desc": ""},
    {"label": "", "desc": "由于 $s_z$ < 0，该表达式保证支座高度遵循物理规律正确减小"},
    {"label": "深远影响", "desc": ""},
    {"label": "P-Δ力矩", "desc": "：势能泛函全微分时自然涌现 ∂/∂θ(P·(h₀+s_z)·sinθ) → 受压倾斜支座重力倾覆加剧"},
    {"label": "", "desc": "若未修正，系统会荒谬地表现出「轴压越大越自动复位」的自稳定假象——极其危险！"},
], formulas=[{"latex": r"h_{eff} = $h_0$ + $s_z$ \; (\text{NOT } $h_0$ - $s_z$), \quad M_{P-\Delta} = P \cdot h_{eff} \cdot \sin\theta", "label": "高度表达式修正"}])

# ---- SLIDE 24: Nonlinear stiffness ----
make_content_slide("支座非线性刚度表达式", [
    {"label": "水平剪切刚度 $K_h(s)$", "desc": ""},
    {"label": "", "desc": "$K_h(s)$=K_h0·(1-tanh(α·|s|/$T_r$)) —— 双曲正切函数描述后屈曲软化"},
    {"label": "切线导数", "desc": "：$dK_h/ds$ = -K_h0·α/$T_r$·sech²(α·|s|/$T_r$) —— Newton-Raphson解析雅可比"},
    {"label": "竖向压缩刚度 $K_v$", "desc": ""},
    {"label": "", "desc": "$K_v$ = $E_c$·A/$T_r$ (恒定线弹性) —— $E_c$为修正压缩模量，一阶导数为0"},
    {"label": "旋转抗弯刚度 $K_\\theta$(s)", "desc": ""},
    {"label": "", "desc": "$K_\\theta$(s) = 0.5·$K_\\theta$0·(1-$C_\\theta$·|s|) —— 基于上下对称均分+线性退化假设"},
    {"label": "", "desc": "d$K_\\theta$/ds = -0.5·$K_\\theta$0·$C_\\theta$ —— 常数导数，简化的线性退化特征"},
], formulas=[
    {"latex": r"$K_h(s)$ = $K_{h0}$ \cdot (1 - \tanh(\alpha \cdot |s|/$T_r$)), \quad $K_v$ = \frac{$E_c$ \cdot A}{$T_r$}", "label": "水平+竖向刚度"},
    {"latex": r"K_\theta(s) = \frac{1}{2} K_{\theta 0} \cdot (1 - C_\theta \cdot |s|), \quad \frac{dK_\theta}{ds} = -\frac{1}{2} K_{\theta 0} C_\theta", "label": "旋转刚度"},
])

# ---- SLIDE 25: Single isolator 4-DOF ----
make_content_slide("单水平隔震支座：4-DOF数学模型", [
    {"label": "内部自由度 (微观)", "desc": "：v = [s_x, $s_z$, θ_h, θ_v]^T —— 纯剪、竖压、形心转角、顶板转角"},
    {"label": "外部输出 (宏观)", "desc": "：F_local = [$V_b$, $P_b$, $M_b$]^T —— 剪力、轴力、弯矩"},
    {"label": "大变形几何映射 T(v)", "desc": "：将微观内部变形非线性映射为宏观顶板位移"},
    {"label": "4个残差方程", "desc": "：R₁水平平衡 · R₂竖向平衡 · R₃倾覆平衡 · R₄顶板平衡"},
    {"label": "刚度凝聚", "desc": ""},
    {"label": "", "desc": "K_local = (T₁ᵀ·K_eq⁻¹·T₁)⁻¹ —— 4×4→3×3宏观刚度"},
], formulas=[
    {"latex": r"v = [s_x, $s_z$, \theta_h, \theta_v]^T, \quad $F_{local}$ = [$V_b$, $P_b$, $M_b$]^T", "label": "4-DOF定义"},
    {"latex": r"K_{local} = \left( T_1^T \cdot $K_{eq}$^{-1} \cdot T_1 \right)^{-1} \in \mathbb{R}^{3\times 3}", "label": "顶板刚度凝聚"},
])

# ---- SLIDE 26: Naming transition ----
make_content_slide("刚度命名体系：单支座→串联的统一过渡", [
    {"label": "命名冲突", "desc": ""},
    {"label": "单支座", "desc": "：$K_\\theta$h (底部/horizontal)、$K_\\theta$v (顶部/vertical-top)"},
    {"label": "串联", "desc": "：存在两个支座+中间板 → 继续用 $K_\\theta$h/$K_\\theta$v 导致下标冲突与物理混淆"},
    {"label": "统一过渡规则", "desc": "：基于位置属性命名"},
    {"label": "$K_{bot}$", "desc": "：Bottom - 底部弹簧刚度"},
    {"label": "$K_{top}$", "desc": "：Top - 顶部弹簧刚度"},
    {"label": "串联实例", "desc": ""},
    {"label": "上支座(#1)", "desc": "：$K_{bot}$₁连中间板、$K_{top}$₁连刚性底板"},
    {"label": "下支座(#2)", "desc": "：$K_{bot}$₂连基础、$K_{top}$₂连中间板"},
    {"label": "中间板枢纽", "desc": "：$K_{top}$₂ + $K_{bot}$₁ 共同约束 Θ_mid，这是 R₉ 残差的物理基础"},
], formulas=[{"latex": r"\text{单支座: }K_{\theta h},\;K_{\theta v} \;\xrightarrow{\text{串联}}\; K_{bot},\;K_{top}", "label": "命名过渡"}])

# ---- SLIDE 27: Series 8-DOF ----
make_content_slide("串联隔震支座：8-DOF高维强非线性系统", [
    {"label": "8-DOF内部变形向量", "desc": "：v=[s_x1,s_z1,θ_h1,θ_v1,$X_{mid}$,$Z_{mid}$,Θ_mid,s_x2,s_z2,θ_h2,θ_v2]^T"},
    {"label": "全系统几何映射", "desc": "：上下支座各自独立大变形映射 + 中间板运动学协调"},
    {"label": "硬接触惩罚机制", "desc": "：中间板转角超限 → 罚函数法强非线性接触力矩锁止"},
    {"label": "外部输出", "desc": "：F_local=[$V_b$, $P_b$, $M_b$]^T —— 与单支座完全统一接口"},
    {"label": "凝聚输出", "desc": "：从8-DOF→3-DOF宏观反力，flag参数自动切换 'single'/'series'"},
], formulas=[{"latex": r"v_{series} = [$s_{x1}$, $s_{z1}$, \theta_{h1}, \theta_{v1}, X_{mid}, Z_{mid}, \Theta_{mid}, $s_{x2}$, $s_{z2}$, \theta_{h2}, \theta_{v2}]^T", "label": "串联8-DOF"},
    {"latex": r"$K_{global}$^{series} = \left( T_{series} \cdot K_{int,8\times 8}^{-1} \cdot T_{series}^T \right)^{-1}", "label": "串联凝聚"},
])

# ---- SLIDE 28: Schur complement verification ----
make_content_slide("K_global 公式的数学严密性校核", [
    {"label": "核心问题", "desc": "：K_global = (T·K_int⁻¹·Tᵀ)⁻¹ 与舒尔补的等价性在数学上是否严密？"},
    {"label": "与传统静力凝聚的区别", "desc": ""},
    {"label": "", "desc": "$U_{ext}$ 不是 v_int 的简单子集——存在高度非线性几何映射 T(v)"},
    {"label": "", "desc": "无法对 K 做简单的 Partition 切块操作"},
    {"label": "严密推导 (混合列式 → 鞍点系统)", "desc": ""},
    {"label": "", "desc": "增量矩阵：[K_int  Tᵀ; T  0]·[dv; dF] = [0; dU]"},
    {"label": "", "desc": "对左上角 K_int 求舒尔补消去 dv → 得宏观刚度"},
    {"label": "校核结论", "desc": "：公式完全正确。T·K_int⁻¹·Tᵀ = 系统对外宏观切线柔度 (Compliance)，求逆即为有效宏观刚度"},
], formulas=[{"latex": r"\begin{bmatrix} $K_{int}$ & T^T \ T & 0 \end{bmatrix} \begin{bmatrix} dv \ dF \end{bmatrix} = \begin{bmatrix} 0 \ dU \end{bmatrix} \Rightarrow $K_{global}$ = (T $K_{int}$^{-1} T^T)^{-1}", "label": "鞍点系统推导"}])

# ---- SLIDE 29: Superstructure ----
make_content_slide("上部多质点体系与全局耦合机理", [
    {"label": "质量矩阵", "desc": "：绝对坐标系下 M = diag(m₁, m₂, ..., m_N) —— 完美对角矩阵"},
    {"label": "弹性刚度", "desc": "：采用D值法组装层间刚度矩阵"},
    {"label": "动态 P-Δ 效应", "desc": ""},
    {"label": "层间侧向刚度", "desc": "：随楼层轴力实时修正，底部层间软化最为显著"},
    {"label": "底板转动刚度", "desc": "：整体削弱 → 宏观倾覆效应通过实时 $K_G$ 矩阵修正"},
    {"label": "几何映射 T_j", "desc": "：第j个支座位置 → 局部坐标系→底板质心全局坐标系"},
    {"label": "四步闭环工作流", "desc": "：下发位移 → 微观内平衡 → 静力凝聚Pull-back → 全局组装 → 时间步进"},
], formulas=[{"latex": r"$K_{global}$^{struct} = $K_{elastic}$ + $K_G$(P), \quad $K_G$ = f(P_i, h_i, \Delta u_i)", "label": "上部P-Δ刚度"}])

# ---- SLIDE 30: Workflow summary ----
make_content_slide("全局数值求解闭环工作流", [
    {"label": "Step 1 下发试探位移", "desc": "：主控系统基于上一时间步结果，利用 Newmark-β 法下发 $\Delta U_{top}$"},
    {"label": "Step 2 微观内平衡", "desc": "：支座内部 Newton-Raphson 求解 R(v)=0$ → 得 v* 和局部 K_int"},
    {"label": "Step 3 静力凝聚 Pull-back", "desc": "：利用映射矩阵 T_j^T 将局部 K_int、F_int 拉回到底板质心"},
    {"label": "Step 4 全局组装", "desc": "：合成全局有效刚度-质量系统 → 推进外层主结构方程收敛"},
    {"label": "自适应时间步长", "desc": "：SmartAnalyze — 检测收敛失败自动细分步长(min_dt=1e-6, tol=1e-7)"},
    {"label": "关键保障", "desc": "：双重迭代架构 + 舒尔补凝聚 → 非线性突变下稳定高精度数值解"},
])

# ---- SLIDE 31: Summary ----
s = add_blank_slide(DARK_BG)
shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.8), Inches(1.5), Inches(0.06))
shp.fill.solid(); shp.fill.fore_color.rgb = TEAL; shp.line.fill.background()
txBox = s.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(8.5), Inches(0.6))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "总结与展望"
p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = FONT_TITLE

points = [
    ("理论完备性", "3DOF→5DOF→9DOF演进基于物理对称性与变分原理，数学推导严密，向下兼容验证通过"),
    ("可计算性", "舒尔补刚度凝聚 + 嵌套Newton-Raphson → 控制全局规模的同时保证精度"),
    ("接口统一", "单/串联支座统一flag接口，相同3-DOF宏观输出 → 代码复用性高"),
    ("已知局限", "几何刚度截断缺陷 · 割线退化模型替代Bouc-Wen滞回 · 二维地震验证待补"),
    ("后续方向", "完善$K_{geo}$高阶项 · Bouc-Wen循环滞回模型 · 三维双向地震动验证 · 碰撞参数标定"),
]
for i, (t, d) in enumerate(points):
    yy = Inches(1.9 + i * 0.7)
    s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), yy, Inches(0.06), Inches(0.4)).fill.solid()
    s.shapes[0].fill.fore_color.rgb = TEAL  # will be overridden per shape, need separate fill
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), yy + Inches(0.05), Inches(0.05), Inches(0.35))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
    txBox = s.shapes.add_textbox(Inches(1.0), yy, Inches(2.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = t
    p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = ICE_BLUE; p.font.name = FONT_TITLE
    txBox2 = s.shapes.add_textbox(Inches(3.5), yy, Inches(5.5), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = d
    p2.font.size = Pt(11); p2.font.color.rgb = LIGHT_GRAY; p2.font.name = FONT_BODY

# Bottom
shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(5.1), Inches(1.5), Inches(0.04))
shp.fill.solid(); shp.fill.fore_color.rgb = CORAL; shp.line.fill.background()

# ---- SLIDE 32: Thank you ----
s = add_blank_slide(DARK_BG)
txBox = s.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8.6), Inches(1.0))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "谢谢！"
p.font.size = Pt(42); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = FONT_TITLE
p2 = tf.add_paragraph()
p2.text = "串联支座隔震结构 · 9自由度理论框架审查"
p2.font.size = Pt(16); p2.font.color.rgb = TEAL; p2.font.name = FONT_BODY
shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(3.1), Inches(2.0), Inches(0.05))
shp.fill.solid(); shp.fill.fore_color.rgb = CORAL; shp.line.fill.background()
txBox2 = s.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(8.6), Inches(0.4))
tf2 = txBox2.text_frame
p3 = tf2.paragraphs[0]
p3.text = "iFLOW 结构抗震分析团队  |  2026年5月"
p3.font.size = Pt(12); p3.font.color.rgb = MUTED; p3.font.name = FONT_BODY

# ============================================================
# SAVE
# ============================================================
output_path = "C:/Users/Fqdsooner/Desktop/iFLOW/ClaudeCodeEnvironmentSupporter/串联支座9自由度模型审查报告_v3.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
